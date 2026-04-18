import csv
import io
import os
import re
import json
import html
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from typing import List, Dict, Any, Optional, Tuple
from fastapi import UploadFile


SUPPORTED_FORMATS = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls": "application/vnd.ms-excel",
    ".xmind": "application/xmind",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".md": "text/markdown",
    ".txt": "text/plain",
    ".csv": "text/csv",
    ".html": "text/html",
    ".htm": "text/html",
    ".json": "application/json",
    ".yaml": "application/x-yaml",
    ".yml": "application/x-yaml",
    ".xml": "application/xml",
    ".epub": "application/epub+zip",
}


def detect_format(filename: str) -> Optional[str]:
    ext = os.path.splitext(filename)[1].lower()
    return ext if ext in SUPPORTED_FORMATS else None


async def parse_file(upload_file: UploadFile) -> Tuple[str, str]:
    ext = detect_format(upload_file.filename)
    if not ext:
        raise ValueError(f"Unsupported file format: {upload_file.filename}")

    content = await upload_file.read()

    if ext == ".pdf":
        return parse_pdf(content, upload_file.filename)
    elif ext in (".docx", ".doc"):
        return parse_word(content, upload_file.filename)
    elif ext in (".xlsx", ".xls"):
        return parse_excel(content, upload_file.filename)
    elif ext == ".xmind":
        return parse_xmind(content, upload_file.filename)
    elif ext == ".pptx":
        return parse_pptx(content, upload_file.filename)
    elif ext in (".md", ".txt"):
        return content.decode("utf-8", errors="ignore"), upload_file.filename
    elif ext == ".csv":
        return parse_csv(content, upload_file.filename)
    elif ext in (".html", ".htm"):
        return parse_html(content, upload_file.filename)
    elif ext == ".json":
        return parse_json(content, upload_file.filename)
    elif ext in (".yaml", ".yml"):
        return parse_yaml(content, upload_file.filename)
    elif ext == ".xml":
        return parse_xml(content, upload_file.filename)
    elif ext == ".epub":
        return parse_epub(content, upload_file.filename)
    else:
        raise ValueError(f"Unsupported format: {ext}")


def parse_pdf(content: bytes, filename: str) -> Tuple[str, str]:
    import fitz

    doc = fitz.open(stream=content, filetype="pdf")
    texts = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            texts.append(f"--- Page {page_num + 1} ---\n{text}")
    doc.close()
    return "\n\n".join(texts), filename


def parse_word(content: bytes, filename: str) -> Tuple[str, str]:
    from docx import Document

    doc = Document(io.BytesIO(content))
    texts = []
    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                texts.append(row_text)
    return "\n\n".join(texts), filename


def parse_excel(content: bytes, filename: str) -> Tuple[str, str]:
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".xls":
        import xlrd

        wb = xlrd.open_workbook(file_contents=content)
        texts = []
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            texts.append(f"--- Sheet: {sheet_name} ---")
            for row_idx in range(ws.nrows):
                row = ws.row_values(row_idx)
                row_text = " | ".join(str(cell) if cell else "" for cell in row)
                if row_text.strip():
                    texts.append(row_text)
        return "\n\n".join(texts), filename

    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    texts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        texts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                texts.append(row_text)
    wb.close()
    return "\n\n".join(texts), filename


def parse_xmind(content: bytes, filename: str) -> Tuple[str, str]:
    import zipfile
    import xml.etree.ElementTree as ET

    texts = []
    with zipfile.ZipFile(io.BytesIO(content)) as zf:
        if "content.json" in zf.namelist():
            import json as _json

            data = _json.loads(zf.read("content.json"))
            texts.extend(_extract_xmind_json(data))
        elif "content.xml" in zf.namelist():
            xml_content = zf.read("content.xml")
            root = ET.fromstring(xml_content)
            texts.extend(_extract_xmind_xml(root))
    return "\n\n".join(texts), filename


def _extract_xmind_json(data: Any, prefix: str = "") -> List[str]:
    texts = []
    if isinstance(data, dict):
        title = data.get("title", "")
        if title:
            full_title = f"{prefix} > {title}" if prefix else title
            texts.append(full_title)
            for child in data.get("children", {}).get("attached", []):
                texts.extend(_extract_xmind_json(child, full_title))
    elif isinstance(data, list):
        for item in data:
            texts.extend(_extract_xmind_json(item, prefix))
    return texts


def _extract_xmind_xml(root: Any, prefix: str = "") -> List[str]:
    texts = []
    for sheet in root.findall(".//{*}sheet"):
        title_elem = sheet.find("{*}title")
        sheet_title = title_elem.text if title_elem is not None else ""
        for topic in sheet.findall(".//{*}topic"):
            title_elem = topic.find("{*}title")
            if title_elem is not None and title_elem.text:
                texts.append(f"{sheet_title} > {title_elem.text}")
    return texts


def parse_pptx(content: bytes, filename: str) -> Tuple[str, str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text)
        if slide_texts:
            texts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(texts), filename


def parse_csv(content: bytes, filename: str) -> Tuple[str, str]:
    text = content.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return "", filename
    header = rows[0]
    texts = []
    for i, row in enumerate(rows[1:], 1):
        row_parts = []
        for j, cell in enumerate(row):
            if j < len(header):
                row_parts.append(f"{header[j]}: {cell.strip()}")
            else:
                row_parts.append(cell.strip())
        texts.append(f"--- Row {i} ---\n" + "\n".join(row_parts))
    return "\n\n".join(texts), filename


def parse_html(content: bytes, filename: str) -> Tuple[str, str]:
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(content, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
    except ImportError:
        raw = content.decode("utf-8", errors="ignore")
        raw = re.sub(
            r"<script[^>]*>.*?</script>", "", raw, flags=re.DOTALL | re.IGNORECASE
        )
        raw = re.sub(
            r"<style[^>]*>.*?</style>", "", raw, flags=re.DOTALL | re.IGNORECASE
        )
        raw = re.sub(r"<[^>]+>", " ", raw)
        raw = html.unescape(raw)
        text = raw
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines), filename


def parse_json(content: bytes, filename: str) -> Tuple[str, str]:
    raw = content.decode("utf-8", errors="ignore")
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        return raw, filename

    texts = []

    def _extract(obj, prefix=""):
        if isinstance(obj, dict):
            for key, value in obj.items():
                new_prefix = f"{prefix}.{key}" if prefix else key
                _extract(value, new_prefix)
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                new_prefix = f"{prefix}[{i}]"
                _extract(item, new_prefix)
        elif isinstance(obj, (str, int, float, bool)):
            texts.append(f"{prefix}: {obj}")

    _extract(data)
    return "\n".join(texts), filename


def parse_yaml(content: bytes, filename: str) -> Tuple[str, str]:
    try:
        import yaml

        raw = content.decode("utf-8", errors="ignore")
        data = yaml.safe_load(raw)
    except ImportError:
        return content.decode("utf-8", errors="ignore"), filename

    texts = []

    def _extract(obj, prefix=""):
        if isinstance(obj, dict):
            for key, value in obj.items():
                new_prefix = f"{prefix}.{key}" if prefix else str(key)
                _extract(value, new_prefix)
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                new_prefix = f"{prefix}[{i}]"
                _extract(item, new_prefix)
        elif isinstance(obj, (str, int, float, bool)):
            texts.append(f"{prefix}: {obj}")

    _extract(data)
    return "\n".join(texts), filename


def parse_xml(content: bytes, filename: str) -> Tuple[str, str]:
    raw = content.decode("utf-8", errors="ignore")
    root = ET.fromstring(raw)
    texts = []

    def _extract(elem, prefix=""):
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        text = (elem.text or "").strip()
        attribs = ", ".join(f"{k}={v}" for k, v in elem.attrib.items())
        if text:
            full_tag = f"{prefix}/{tag}" if prefix else tag
            if attribs:
                texts.append(f"{full_tag} ({attribs}): {text}")
            else:
                texts.append(f"{full_tag}: {text}")
        for child in elem:
            new_prefix = f"{prefix}/{tag}" if prefix else tag
            _extract(child, new_prefix)

    _extract(root)
    return "\n".join(texts), filename


def parse_epub(content: bytes, filename: str) -> Tuple[str, str]:
    texts = []
    with zipfile.ZipFile(io.BytesIO(content)) as zf:
        namelist = zf.namelist()
        spine_order = []
        try:
            if "META-INF/container.xml" in namelist:
                container = ET.fromstring(zf.read("META-INF/container.xml"))
                opf_path = container.find(".//{*}rootfile").get(
                    "full-path", "OEBPS/content.opf"
                )
            else:
                opf_path = "OEBPS/content.opf"

            opf_dir = os.path.dirname(opf_path)
            if opf_path in namelist:
                opf_root = ET.fromstring(zf.read(opf_path))
                ns = {"opf": "http://www.idpf.org/2007/opf"}
                manifest = {}
                for item in opf_root.findall(".//{*}item"):
                    item_id = item.get("id", "")
                    item_href = item.get("href", "")
                    manifest[item_id] = item_href

                spine = opf_root.find(".//{*}spine")
                if spine is not None:
                    for itemref in spine.findall("{*}itemref"):
                        idref = itemref.get("idref", "")
                        if idref in manifest:
                            spine_order.append(manifest[idref])

            for rel_path in spine_order:
                full_path = os.path.join(opf_dir, rel_path).replace("\\", "/")
                if full_path in namelist:
                    xhtml = zf.read(full_path).decode("utf-8", errors="ignore")
                    try:
                        from bs4 import BeautifulSoup

                        soup = BeautifulSoup(xhtml, "html.parser")
                        for tag in soup(["script", "style"]):
                            tag.decompose()
                        chapter_text = soup.get_text(separator="\n")
                    except ImportError:
                        chapter_text = re.sub(r"<[^>]+>", " ", xhtml)
                        chapter_text = html.unescape(chapter_text)
                    lines = [
                        line.strip()
                        for line in chapter_text.splitlines()
                        if line.strip()
                    ]
                    if lines:
                        texts.append(
                            f"--- Chapter: {rel_path} ---\n" + "\n".join(lines)
                        )

        except Exception:
            for name in namelist:
                if name.endswith((".xhtml", ".html", ".htm")):
                    raw = zf.read(name).decode("utf-8", errors="ignore")
                    clean = re.sub(r"<[^>]+>", " ", raw)
                    texts.append(html.unescape(clean))

    return "\n\n".join(texts), filename


def auto_detect_doc_type(text: str, filename: str) -> str:
    return "document"


def auto_detect_module(text: str, filename: str) -> str:
    lower_name = filename.lower()
    name_no_ext = os.path.splitext(lower_name)[0]
    module_patterns = [
        r"(用户中心|用户管理|user[\s_-]?center|account)",
        r"(支付|payment|pay)",
        r"(登录|login|auth)",
        r"(订单|order)",
        r"(商品|product|goods)",
        r"(消息|message|notification)",
        r"(权限|permission|role|rbac)",
        r"(配置|config|setting)",
        r"(报表|report|dashboard)",
        r"(审批|approval|workflow)",
    ]
    for pattern in module_patterns:
        match = re.search(pattern, lower_name, re.IGNORECASE)
        if match:
            return match.group(1)
    return "通用模块"
